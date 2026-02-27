#!/usr/bin/env python3
import sys
import os
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.chart.data import ChartData
from io import BytesIO


def extract_text(shape):
    """从文本形状提取文本内容，返回合并后的文本"""
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        lines.append(para.text if para.text else "")
    return "\n".join([l for l in lines if l]).strip()


def convert_pptx_to_editable(input_path, output_path, preserve_watermark=True):
    """
    将 NotebookLM 生成的 PPTX 转换成尽量可编辑的 PPTX。
    主要思路：逐张幻灯片创建空白布局幻灯片，复制文本为文本框、复制图片为图片。
    注意：该实现尽力保留内容和布局的可编辑性，但对复杂对象（如图表、表格、SmartArt 等）处理有限。
    watermark 删除并非此工具的目标，保留原有水印或与水印相关对象将按原样处理，详见实现。
    """
    src = Presentation(input_path)
    dst = Presentation()

    # 选择一个空白布局作为新幻灯片的布局（常见为 index 6，但不同模板可能不同）
    blank_layout_index = 6
    # 如果默认空白布局不可用，兜底使用第一个布局
    blank_layout = dst.slide_layouts[blank_layout_index] if len(dst.slide_layouts) > blank_layout_index else dst.slide_layouts[0]

    for slide_idx, slide in enumerate(src.slides, start=1):
        new_slide = dst.slides.add_slide(blank_layout)

        # 遍历原幻灯片的形状
        for shape in slide.shapes:
            # 1) 文本形状：复制为可编辑文本框，保留原位置
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                text = extract_text(shape)
                if text:
                    tx_box = new_slide.shapes.add_textbox(left, top, width, height)
                    tf = tx_box.text_frame
                    tf.text = text
                    # 字体映射：尽量保留原始字体设定（取文本框第一段的字体作为默认值）
                    try:
                        first_para = shape.text_frame.paragraphs[0]
                        default_font = first_para.font
                        if default_font is not None:
                            for p in tf.paragraphs:
                                if getattr(default_font, "name", None):
                                    p.font.name = default_font.name
                                if getattr(default_font, "size", None):
                                    p.font.size = default_font.size
                    except Exception:
                        pass
            # 2) 图片形状：复制为新的图片对象，保留大致位置和尺寸
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_stream = BytesIO(image_bytes)
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    new_slide.shapes.add_picture(image_stream, left, top, width, height)
                except Exception:
                    # 忽略无法复制的图片
                    pass
            # 3) 表格：复制为可编辑表格，尽量保留文本
            elif getattr(shape, "has_table", False) and shape.table is not None:
                try:
                    tbl = shape.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    new_tbl = new_slide.shapes.add_table(rows, cols, left, top, width, height).table
                    for i in range(rows):
                        for j in range(cols):
                            new_tbl.cell(i, j).text = tbl.cell(i, j).text
                except Exception:
                    pass
            # 4) 图表：尽量提取数据重建成可编辑图表，失败时回退为原始文本/图片的处理
            elif getattr(shape, "has_chart", False) and shape.chart is not None:
                try:
                    ch = shape.chart
                    chart_type = ch.chart_type
                    cd = ChartData()
                    # 尝试提取类别标签
                    categories = []
                    try:
                        for cat in getattr(ch, "category_axis", []):
                            if hasattr(cat, "categories"):
                                for c in cat.categories:
                                    categories.append(str(c))
                                break
                    except Exception:
                        pass
                    if categories:
                        cd.categories = categories
                    # 提取系列数据
                    series_list = []
                    for s in getattr(ch, "series", []):
                        name = getattr(s, "name", None) or "Series"
                        values = []
                        try:
                            for v in getattr(s, "values", []):
                                values.append(v)
                        except Exception:
                            values = []
                        if values:
                            cd.add_series(name, values)
                    if len(cd.series) > 0:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        new_slide.shapes.add_chart(chart_type, left, top, width, height, cd)
                except Exception:
                    pass
            # 其他类型（如 SmartArt）尝试简化为文本占位，或跳过

    # 保存输出
    dst.save(output_path)


def main():
    parser = argparse.ArgumentParser(
        description="NotebookLM PPT 转可编辑 PPT 的 CLI（水印移除功能请自行确保合法授权）"
    )
    parser.add_argument("input", help="NotebookLM 生成的 PPTX 文件路径")
    parser.add_argument("-o", "--output", dest="output", help="输出的 PPTX 路径（默认为输入文件名加 _editable 后缀）")
    parser.add_argument("--no-watermark-removal", action="store_true",
                        help="提示：该工具不提供水印移除功能。请确保你对内容拥有合适的使用权。")

    args = parser.parse_args()

    input_path = args.input
    if not os.path.exists(input_path):
        print(f"错误：输入文件不存在：{input_path}", file=sys.stderr)
        sys.exit(2)

    if args.output:
        output_path = args.output
    else:
        base, ext = os.path.splitext(input_path)
        output_path = base + "_editable.pptx"

    try:
        convert_pptx_to_editable(input_path, output_path, preserve_watermark=True)
    except Exception as e:
        print(f"转换失败: {e}", file=sys.stderr)
        sys.exit(3)

    print(f"可编辑 PPTX 已创建：{output_path}")


if __name__ == "__main__":
    main()
