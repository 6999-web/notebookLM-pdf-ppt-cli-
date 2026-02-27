import os
import tempfile
import subprocess
from pptx import Presentation
from pptx.util import Inches


def _create_test_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)
    # Add text box
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = tx.text_frame
    tf.text = "Hello from NotebookLM"
    # Add image
    from PIL import Image, ImageDraw
    img_path = path.replace('.pptx', '_img.png')
    img = Image.new('RGB', (200, 100), color=(73, 109, 137))
    d = ImageDraw.Draw(img)
    d.text((10, 40), "Test Image", fill=(255, 255, 0))
    img.save(img_path)
    slide.shapes.add_picture(img_path, Inches(1), Inches(3), Inches(2), Inches(1.5))
    prs.save(path)


def test_basic_text_and_image_conversion():
    with tempfile.TemporaryDirectory() as tmp:
        input_pptx = os.path.join(tmp, 'input.pptx')
        output_pptx = os.path.join(tmp, 'output_editable.pptx')
        _create_test_pptx(input_pptx)
        # Run the CLI converter
        subprocess.run(['python', '-m', 'notebook_editor.main', input_pptx, '-o', output_pptx], check=True)
        # Validate output
        out = Presentation(output_pptx)
        assert len(out.slides) >= 1
        found_text = False
        found_image = False
        for s in out.slides:
            for sh in s.shapes:
                if getattr(sh, 'has_text_frame', False) and sh.text_frame and any(p.text for p in sh.text_frame.paragraphs):
                    text = "%s" % sh.text_frame.text
                    if "Hello from NotebookLM" in text:
                        found_text = True
                if sh.shape_type.name == 'PICTURE' or getattr(sh, 'image', None) is not None:
                    found_image = True
        assert found_text
        assert found_image
