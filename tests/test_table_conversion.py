import os
import tempfile
import subprocess
from pptx import Presentation
from pptx.util import Inches


def _create_table_pptx(path):
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)
    rows, cols = 2, 2
    left, top, width, height = Inches(1), Inches(1), Inches(6), Inches(1.5)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    tbl_shape.cell(0, 0).text = 'A1'
    tbl_shape.cell(0, 1).text = 'A2'
    tbl_shape.cell(1, 0).text = 'B1'
    tbl_shape.cell(1, 1).text = 'B2'
    prs.save(path)


def test_table_conversion_preserves_text():
    with tempfile.TemporaryDirectory() as tmp:
        input_pptx = os.path.join(tmp, 'input_table.pptx')
        output_pptx = os.path.join(tmp, 'output_table_editable.pptx')
        _create_table_pptx(input_pptx)
        subprocess.run(['python', '-m', 'notebook_editor.main', input_pptx, '-o', output_pptx], check=True)
        out = Presentation(output_pptx)
        assert len(out.slides) >= 1
        found_table = False
        for s in out.slides:
            for sh in s.shapes:
                if getattr(sh, 'has_table', False) and sh.table is not None:
                    tbl = sh.table
                    assert tbl.cell(0, 0).text == 'A1'
                    assert tbl.cell(0, 1).text == 'A2'
                    assert tbl.cell(1, 0).text == 'B1'
                    assert tbl.cell(1, 1).text == 'B2'
                    found_table = True
        assert found_table
