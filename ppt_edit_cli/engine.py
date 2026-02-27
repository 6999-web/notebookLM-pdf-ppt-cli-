import cv2
import numpy as np
import os
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import google.generativeai as genai
from skimage.restoration import inpaint
from .utils import px_to_emu, hex_to_rgb

class PPTEngine:
    def __init__(self, api_key=None):
        # Initialize PaddleOCR
        # Initializing here will download models on first run
        self.ocr = PaddleOCR(use_angle_cls=True, lang='ch', show_log=False)
        
        # Initialize LLM (Gemini)
        if api_key:
            genai.configure(api_key=api_key)
            self.model = genai.GenerativeModel('gemini-1.5-flash') # Using flash for speed
        else:
            self.model = None

    def process_image(self, img_path):
        """Main processing pipeline for a single image."""
        img = cv2.imread(img_path)
        if img is None:
            raise ValueError(f"Could not read image: {img_path}")
        h, w, _ = img.shape
        
        # 1. OCR and Element Detection
        ocr_results = self.ocr.ocr(img_path, cls=True)
        
        # 2. Semantic Reconstruction (LLM)
        structured_content = self._reconstruct_semantics(ocr_results, w, h)
        
        # 3. Background Extraction & Inpainting (Clean the slide)
        # We use high-confidence text boxes and image boxes to inpaint the original background
        background_path = self._extract_background(img, img_path, structured_content)
        
        return {
            'width': w,
            'height': h,
            'background': background_path,
            'elements': structured_content
        }

    def _extract_background(self, img, original_path, structured_content, remove_watermark=True):
        """Identify and extract background with inpainting of text regions and watermarks."""
        bg_path = original_path.replace('.', '_bg.')
        
        # Create mask for regions to inpaint
        mask = np.zeros(img.shape[:2], dtype=bool)
        for element in structured_content:
            if element.get('type') != 'Image':
                x, y, w, h = element['box']
                pad = 2
                mask[max(0, int(y-pad)):min(img.shape[0], int(y+h+pad)), 
                     max(0, int(x-pad)):min(img.shape[1], int(x+w+pad))] = True
        
        # NotebookLM Watermark (usually bottom right)
        if remove_watermark:
            h, w = img.shape[:2]
            # NotebookLM watermark is typically in the bottom right corner
            # Values from NotebookLM2PPT: r1,r2,c1,c2 = 1536,1598,2627,2863 for a 2867x1600 image
            # We scale these to the current image size
            wm_r1 = int(1536 * h / 1600)
            wm_r2 = int(1598 * h / 1600)
            wm_c1 = int(2627 * w / 2867)
            wm_c2 = int(2863 * w / 2867)
            mask[wm_r1:wm_r2, wm_c1:wm_c2] = True

        if np.any(mask):
            try:
                # Use biharmonic inpainting
                inpainted_img = inpaint.inpaint_biharmonic(img, mask, channel_axis=-1)
                inpainted_img = (inpainted_img * 255).astype(np.uint8)
                cv2.imwrite(bg_path, inpainted_img)
            except Exception as e:
                print(f"Inpainting failed: {e}")
                cv2.imwrite(bg_path, img)
        else:
            cv2.imwrite(bg_path, img) 
            
        return bg_path

    def _reconstruct_semantics(self, ocr_results, width, height):
        """Send OCR data to LLM to group and correct text."""
        if not self.model or not ocr_results or not ocr_results[0]:
            return self._simple_parse(ocr_results, width, height)
            
        ocr_data = []
        for line in ocr_results[0]:
            coords = line[0]
            text, confidence = line[1]
            ocr_data.append({
                "coords": coords,
                "text": text,
                "confidence": confidence
            })
            
        prompt = f"""
        Below is OCR data from a PPT slide (Width: {width}, Height: {height}).
        Each entry has coordinates and text.
        
        Tasks:
        1. Group scattered text fragments into logical blocks:
           - "Title": Main slide heading.
           - "Body": Paragraph or bullet point text.
           - "List": A list of items.
           - "Footer": Page numbers or copyrights.
           - "Table": Data that looks like a table (return as separate rows if possible).
        2. Identify regions that look like IMAGES or CHARTS (return with type "Image" and approximate box).
        3. Correct typos and formatting.
        4. Infer style: font_size (Pt), color (Hex), alignment (left, center, right), bold (boolean).
        5. Detect "Big Artistic Title" (set is_artistic: true).
        
        OCR Data: {ocr_data}
        
        Return a JSON list of objects:
        [{{
          "type": "Title|Body|List|Footer|Table|Image",
          "text": "The content",
          "box": [x, y, w, h],
          "font_size": 24,
          "color": "#000000",
          "align": "left|center|right",
          "bold": false,
          "is_artistic": false
        }}]
        """
        
        try:
            response = self.model.generate_content(prompt)
            text_response = response.text.strip()
            if "```json" in text_response:
                text_response = text_response.split("```json")[1].split("```")[0].strip()
            import json
            return json.loads(text_response)
        except Exception as e:
            print(f"LLM Error: {e}")
            return self._simple_parse(ocr_results, width, height)

    def _simple_parse(self, ocr_results, width, height):
        """Fallback simple parser."""
        blocks = []
        if not ocr_results or not ocr_results[0]: return blocks
        for line in ocr_results[0]:
            coords = line[0]
            text, _ = line[1]
            x_min = min(c[0] for c in coords)
            y_min = min(c[1] for c in coords)
            x_max = max(c[0] for c in coords)
            y_max = max(c[1] for c in coords)
            blocks.append({
                "type": "Body",
                "text": text,
                "box": [x_min, y_min, x_max - x_min, y_max - y_min],
                "font_size": 18,
                "color": "#000000",
                "align": "left",
                "is_artistic": False
            })
        return blocks

    def create_ppt(self, processed_slides, output_path):
        """Generate PPT file from processed data."""
        prs = Presentation()
        # Set slide size for 16:9 (standard EMU values)
        prs.slide_width = 9144000
        prs.slide_height = 5143500 

        for slide_data in processed_slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
            
            # Layer 1: Background
            if slide_data['background']:
                slide.shapes.add_picture(slide_data['background'], 0, 0, prs.slide_width, prs.slide_height)
                
            # Layer 3: Text Blocks
            img_w, img_h = slide_data['width'], slide_data['height']
            for element in slide_data['elements']:
                x, y, w, h = element['box']
                
                # Mapping pixels to PPT slide units
                left = px_to_emu(x, img_w, prs.slide_width / 914400)
                top = px_to_emu(y, img_h, prs.slide_height / 914400) 
                width_emu = px_to_emu(w, img_w, prs.slide_width / 914400)
                height_emu = px_to_emu(h, img_h, prs.slide_height / 914400)
                
                if element.get('type') == 'Image':
                    # If it's identified as an image block, we might want to crop it 
                    # from the original and add it as a separate picture shape for better editability.
                    # For now, we assume it's part of the background.
                    continue

                txBox = slide.shapes.add_textbox(left, top, width_emu, height_emu)
                tf = txBox.text_frame
                tf.text = element.get('text', '')
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.font.size = Pt(element.get('font_size', 18))
                p.font.bold = element.get('bold', False)
                try:
                    p.font.color.rgb = RGBColor(*hex_to_rgb(element.get('color', '#000000')))
                except:
                    pass
                
                # If it's an artistic title, we might make the text nearly transparent or just hidden
                # but for simplicity we keep it editable.
                # if element.get('is_artistic'):
                #     pass 

                align = element.get('align', 'left').lower()
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT

        prs.save(output_path)
